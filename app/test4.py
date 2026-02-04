from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BLUE = RGBColor(0, 51, 102)   # bleu foncé académique
BLACK = RGBColor(0, 0, 0)

prs = Presentation()
title_layout = prs.slide_layouts[0]
content_layout = prs.slide_layouts[1]

# -------- Styles helpers --------
def style_title(shape):
    shape.text_frame.clear()
    p = shape.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = ""
    p.alignment = PP_ALIGN.CENTER

def set_title(slide, text):
    title = slide.shapes.title
    title.text = text
    for p in title.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(40)
            r.font.bold = True
            r.font.color.rgb = BLUE
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def style_body_frame(tf):
    # police, taille, couleur, interligne
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.name = "Calibri"
            r.font.size = Pt(20)
            r.font.color.rgb = BLACK
        p.line_spacing = 1.2

def add_footer(slide, text="Étude – Desoyaux"):
    left = Inches(0.5); top = Inches(6.8); width = Inches(9); height = Inches(0.4)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.RIGHT
    for r in p.runs:
        r.font.size = Pt(11)
        r.font.color.rgb = RGBColor(100, 100, 100)

def add_slide(title_text, bullets):
    slide = prs.slides.add_slide(content_layout)
    # Titre
    slide.shapes.title.text = title_text
    for p in slide.shapes.title.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(32); r.font.bold = True; r.font.color.rgb = BLUE
    # Contenu
    body = slide.placeholders[1].text_frame
    body.clear()
    # 1er point
    p0 = body.paragraphs[0]
    p0.text = bullets[0]; p0.level = 0
    # suivants
    for b in bullets[1:]:
        p = body.add_paragraph(); p.text = b; p.level = 0
    style_body_frame(body)
    add_footer(slide)

# -------- Slide 1: Titre --------
slide = prs.slides.add_slide(title_layout)
set_title(slide, "Desoyaux")
subtitle = slide.placeholders[1]
subtitle.text = "Analyse et solutions de management (version académique)"
for p in subtitle.text_frame.paragraphs:
    for r in p.runs:
        r.font.size = Pt(18); r.font.color.rgb = BLACK
subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
add_footer(slide)

# -------- Slides contenu --------
add_slide("Style de direction", [
    "Participatif et entrepreneurial",
    "Délégation importante",
    "Proximité et réactivité"
])

add_slide("Finalités", [
    "Économiques : croissance et leadership",
    "Sociales : emploi local et intéressement",
    "Environnementales : recyclage et maîtrise des matières"
])

add_slide("Diagnostic stratégique", [
    "Forces : innovation, réseau, industrialisation intégrée",
    "Faiblesses : dépendance main-d’œuvre, coûts des matières",
    "Opportunités : marché en croissance, cocooning",
    "Menaces : critiques écologiques, inflation, géopolitique"
])

add_slide("Stratégie de domaine", [
    "Différenciation par l’innovation",
    "Brevets, R&D, gamme diversifiée"
])

add_slide("Politique d’innovation", [
    "Objectif : conserver l’avantage concurrentiel",
    "Investissements et dépôts de brevets",
    "Diversification de l’offre",
    "Recyclage intégré"
])

add_slide("Problème identifié", [
    "Manque de main-d’œuvre qualifiée pour la pose",
    "Frein au développement"
])

add_slide("Solutions", [
    "Formation et partenariats",
    "Simplification des installations",
    "Attractivité et fidélisation RH",
    "Réseau de sous-traitants certifiés"
])

add_slide("Conclusion", [
    "Leader porté par l’innovation",
    "Défi clé : ressources humaines",
    "Réponse : former, simplifier, attirer, sous-traiter"
])

prs.save("Desoyaux.pptx")
print("✅ Fichier 'Desoyaux.pptx' généré.")

