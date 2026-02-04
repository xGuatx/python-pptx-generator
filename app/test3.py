#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
Script Python utilisant python-pptx pour générer un PowerPoint "ludique mais pro"
contenant 9 diapositives sur le thème "CES 2025: Innovations in AI and Robotics",
avec un minimum de stylisation (barre colorée, texte en puces, etc.).
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# --- Contenu textuel des slides : titre + liste de puces ---
SLIDES_CONTENT = [
    {
        "title": "CES 2025: Innovations in AI and Robotics",
        "theme_color": RGBColor(0x0B, 0x64, 0x92),  # Exemple: bleu/grisé
        "bullets": [
            "Introduction",
            "AI and Robotics in the Spotlight",
            "AI Robots for Social and Emotional Interaction",
            "AI Robots for Autonomy and Physical Assistance",
            "AI in Innovation and Creativity",
            "NVIDIA’s AI and Graphics Advancements",
            "Transition – Bridging AI Robotics and NVIDIA’s Tech",
            "Conclusion",
        ],
    },
    {
        "title": "Introduction – CES 2025 Overview",
        "theme_color": RGBColor(0x13, 0x55, 0x1E),  # Exemple: vert foncé
        "bullets": [
            (
                "Global Stage for Tech Innovation: CES 2025 in Las Vegas drew over "
                "141,000 attendees and 4,500 exhibitors from around the world (ces.tech). "
                "As the premier consumer electronics show, it showcased breakthrough "
                "technologies across AI, robotics, smart homes, automotive, and more. "
                "Industry leaders and startups alike use CES to debut products and set "
                "trends, making it a bellwether event for the tech industry’s future."
            ),
            (
                "AI and Robotics in the Spotlight: A key theme of CES 2025 was the "
                "pervasive presence of artificial intelligence – “AI was everywhere,” "
                "turning once-speculative ideas into practical applications "
                "(rolandberger.com). From intelligent personal robots to autonomous "
                "vehicles, exhibitors demonstrated how advanced AI is enhancing devices’ "
                "capabilities. This year marked a shift toward AI-driven solutions "
                "addressing real-world challenges (e.g. aging populations, personalized "
                "entertainment), underlining CES’s role in shaping how technology meets "
                "human needs."
            ),
        ],
    },
    {
        "title": "AI and Robotics in the Spotlight",
        "theme_color": RGBColor(0x7D, 0x2E, 0x8D),  # Exemple: violet
        "bullets": [
            (
                "This year, AI was everywhere at CES 2025, shifting towards real-world "
                "applications like aging population support and personalized "
                "entertainment. Robots, from personal companions to autonomous vehicles, "
                "highlighted how advanced AI capabilities are permeating consumer "
                "electronics."
            ),
        ],
    },
    {
        "title": "AI Robots for Social and Emotional Interaction",
        "theme_color": RGBColor(0xC4, 0x65, 0x00),  # Exemple: marron/orangé
        "bullets": [
            (
                "Romi – A palm-sized emotional-support AI robot by Mixi (Japan) that "
                "engages in natural, empathic conversations. Designed to combat loneliness "
                "and anxiety, Romi can hold complex real-time dialogue without the user "
                "needing to pause – it even interjects with relevant comments based on "
                "what it “sees” and “hears” (ces.tech). Powered by in-house deep learning "
                "models, it strives to provide genuine empathy and companionship, using a "
                "cute animated face and playful gestures to simulate human-like "
                "interaction (reuters.com)."
            ),
            (
                "TCL “Ai Me” – A modular AI companion robot unveiled as a concept at CES "
                "2025, aimed at supporting children, seniors, and families. This owl-eyed "
                "robot sits in a detachable “space capsule” base and can accompany users "
                "anywhere (even mounting in a smart car) (designboom.com). It has "
                "expressive digital eyes and a friendly voice, using cameras and sensors "
                "to recognize people and surroundings. Ai Me can record family moments, "
                "act as a home security sentry, and control smart home devices – a cute "
                "guardian and assistant fostering warmth and emotional connection "
                "(designboom.com)."
            ),
        ],
    },
    {
        "title": "AI Robots for Autonomy and Physical Assistance",
        "theme_color": RGBColor(0x00, 0x6A, 0x4E),  # Exemple: vert turquoise
        "bullets": [
            (
                "Bosch “Revol” Smart Crib – An AI-powered connected cradle that helps "
                "parents care for infants. It monitors the baby’s vital signs (heart rate, "
                "breathing) using radar sensors and alerts caregivers via an app to any "
                "issues (homecrux.com). It can automatically rock or play lullabies when "
                "the baby cries and even detect changes in air quality or if the baby’s "
                "face is covered, notifying parents for intervention (housedigest.com). "
                "Bosch also designed the crib to be convertible into a height-adjustable "
                "desk as the child grows (housedigest.com)."
            ),
            (
                "Multi-Function Home Robots: CES 2025 showcased robots that serve as "
                "general home assistants. One prototype robot vacuum not only cleans "
                "floors but also acts as an air purifier, patrols the home for security, "
                "and even delivers items (it can literally bring you a sandwich on demand) "
                "(techradar.com). These all-in-one devices illustrate how AI and robotics "
                "are converging to autonomously manage household tasks – from cleaning "
                "and climate control to basic caregiving."
            ),
            (
                "Humanoid Helpers – “Mirokai” Robot: Enchanted Tools (France) debuted an "
                "updated Mirokai humanoid robot – a friendly, child-sized robot on wheels "
                "designed to assist and engage people. The bright orange Mirokai shown at "
                "CES has been deployed in hospitals as an aid for staff and companion for "
                "patients, and it can serve as a concierge in hospitality settings "
                "(interestingengineering.com, facebook.com). With an adorable, "
                "anime-inspired face and the ability to open doors, fetch objects, and "
                "interact socially, Mirokai provides both practical help and emotional "
                "comfort."
            ),
            (
                "Wearable Exoskeletons: AI-driven robotic exoskeleton suits offer new "
                "ways to assist those with mobility challenges. South Korea’s WIM "
                "exoskeleton (by WiRobotics) demonstrated a 78% improvement in walking "
                "speed and endurance among seniors (rockingrobots.com). These wearable "
                "robots use motors and sensors to augment the user’s strength and "
                "balance, hinting at a future where AI-powered exosuits restore mobility "
                "and independence."
            ),
        ],
    },
    {
        "title": "AI in Innovation and Creativity",
        "theme_color": RGBColor(0xAF, 0x1A, 0x1A),  # Exemple: rouge bordeaux
        "bullets": [
            (
                "AI-Generated Music: Startups and musicians at CES 2025 demonstrated AI "
                "systems composing music and soundscapes. Algorithms can generate original "
                "songs in real time based on a listener’s mood or a few sample inputs. "
                "Will.i.am highlighted the potential of AI-generated music, suggesting it "
                "could emotionally move people as deeply as legendary human artists "
                "(yahoo.com). These tools can aid creators by providing instant song "
                "drafts and personalize entertainment with mood-based DJ mixes."
            ),
            (
                "AI-Assisted Fashion Design: Generative AI design tools propose unique "
                "clothing patterns and styles, helping designers push beyond traditional "
                "boundaries (sganalytics.com). Concepts of virtual try-on experiences and "
                "AI stylists that recommend outfits were featured. For consumers, this "
                "means more personalization, and for creators, rapid prototyping of bold "
                "new designs. Fashion in 2025 embraces AI for creativity and efficiency."
            ),
            (
                "Interactive Holograms: Companies showcased 3D holographic displays that "
                "appear free-floating, no AR glasses required. UK-based Hypervsn wowed "
                "attendees with Blade Runner–style holograms – dancers, virtual assistants "
                "hovering mid-air (esquire.com). Future visions include holographic "
                "spokespeople, entertainers, and interactive AI avatars that you can "
                "gesture to. This may become a new medium for communication and art."
            ),
        ],
    },
    {
        "title": "NVIDIA’s AI and Graphics Advancements at CES 2025",
        "theme_color": RGBColor(0x33, 0x49, 0xA2),  # Exemple: bleu plus soutenu
        "bullets": [
            (
                "GeForce RTX 5000 Series Unveiled: Jensen Huang presented next-gen GeForce "
                "RTX 50-series GPUs (Blackwell architecture). The RTX 5090 can do advanced "
                "path-traced lighting and neural rendering, delivering ultra-realistic "
                "graphics (esquire.com). The RTX 5070 at $550 performs on par with the "
                "prior $1,600 RTX 4090 (esquire.com), a 3× value improvement that makes "
                "high-end gaming and AI processing more accessible."
            ),
            (
                "DLSS 4 – Multi-Frame AI Upscaling: NVIDIA introduced DLSS 4 with Multi "
                "Frame Generation, where AI can generate up to 3 additional frames for "
                "every 1 frame rendered (counterpointresearch.com). This yields massively "
                "higher frame rates and smoother motion. Together with improved ray "
                "tracing, AI-optimized shaders, and Reflex 2, DLSS 4 showcases how AI is "
                "revolutionizing real-time graphics."
            ),
            (
                "AI-Powered Robotics Platforms (NVIDIA Isaac & “Cosmos”): NVIDIA also "
                "revealed “Cosmos,” a suite of foundation AI models that generate "
                "photorealistic virtual worlds governed by physics (reuters.com). This "
                "synthetic data accelerates robotics training, reducing real-world trials. "
                "Cosmos is open-license, pairing with the NVIDIA Isaac platform for robot "
                "simulation and development. Such synergy between AI and graphics paves "
                "the way for faster innovation in robotics."
            ),
        ],
    },
    {
        "title": "Transition – Bridging AI Robotics and NVIDIA’s Tech",
        "theme_color": RGBColor(0x53, 0x53, 0x53),  # Exemple: gris
        "bullets": [
            (
                "Synergy of Hardware and AI: The leaps in AI robotics are intertwined with "
                "advances in computing power from companies like NVIDIA. Many smart robots "
                "run on platforms enabled by AI chips and software, meaning improvements "
                "in GPUs and AI models boost robotic capabilities, while robotics’ real-"
                "time demands push chipmakers to innovate specialized hardware "
                "(rolandberger.com)."
            ),
            (
                "AI + Robotics Convergence: CES 2025 blurred the line between consumer "
                "tech and robotics. The same neural networks accelerating gaming also "
                "empower robot vision and speech. NVIDIA’s keynote launched gaming GPUs "
                "and AI training tools for robots in one go (reuters.com). As robots "
                "become common in homes and workplaces, they leverage better chips, "
                "refined algorithms, and cloud services, reinforcing a virtuous cycle of "
                "AI-driven innovation."
            ),
        ],
    },
    {
        "title": "Conclusion – Key Takeaways and Future Outlook",
        "theme_color": RGBColor(0x4B, 0x2D, 0x21),  # Exemple: brun profond
        "bullets": [
            (
                "Major Innovations Recap: CES 2025 showed how AI and robotics are "
                "transforming daily life. Emotional-support robots like Romi and Ai Me, "
                "autonomous devices (smart cribs, exoskeletons), AI-generated music, "
                "AI-assisted fashion, and NVIDIA’s powerful chips/platforms all point to "
                "a more connected, human-centric future."
            ),
            (
                "Benefits to People and Society: Social robots can alleviate loneliness "
                "and provide care for aging populations; assistive robots improve "
                "accessibility and safety; AI creative tools empower artists and "
                "personalize entertainment. Used thoughtfully, these technologies enhance "
                "quality of life, efficiency, and well-being."
            ),
            (
                "Challenges and Considerations: Privacy, reliability, and ethics must be "
                "addressed. Devices collecting personal data (cameras, mics) need secure "
                "handling. Society must ensure safe, trustworthy AI companions, navigate "
                "job displacement, and establish norms/regulations for AI behavior."
            ),
            (
                "Future Vision – “Smart” Everything: From more natural human–AI "
                "interaction to seamless home integration and advanced medical/rehab uses, "
                "the future hinted at by CES 2025 is one of synergy between people and "
                "AI-driven machines. By balancing innovation with responsible "
                "development, we may soon see intelligent robots and AI systems working "
                "hand-in-hand with humans, enriching lives in ways we’re just beginning "
                "to imagine."
            ),
        ],
    },
]


def add_styled_slide(prs, title, bullets, theme_color):
    """
    Crée une diapositive à partir du layout "Blank" et ajoute :
      - Une barre colorée (rectangle) en haut pour le titre
      - Le titre en texte blanc, en gras
      - Les puces en dessous
    """
    # Récupérer la diapo BLANK (souvent index 6 dans PowerPoint par défaut)
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Obtenir la taille de la diapositive
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # -- 1) Ajouter une forme (rectangle) en haut pour le "header" coloré --
    header_height = Inches(1.0)  # Hauteur du bandeau
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=0,
        top=0,
        width=slide_width,
        height=header_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = theme_color
    shape.line.fill.background()  # pas de bordure visible

    # -- 2) Ajouter un TextBox pour le Titre au-dessus de la barre (visuellement "dans" le rectangle) --
    # On va légèrement décaler le texte à l'intérieur (left=Inches(0.3), top=Inches(0.2), etc.)
    title_box = slide.shapes.add_textbox(
        left=Inches(0.3),
        top=Inches(0.2),
        width=slide_width - Inches(0.6),
        height=header_height - Inches(0.4)
    )
    tf = title_box.text_frame
    tf.text = ""  # on va créer manuellement un paragraphe

    p = tf.add_paragraph()
    p.text = title
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)  # blanc

    # -- 3) Ajouter un TextBox pour le contenu en puces, en dessous du header --
    content_box = slide.shapes.add_textbox(
        left=Inches(0.5),
        top=header_height + Inches(0.3),  # on laisse un petit espace sous la barre
        width=slide_width - Inches(1.0),
        height=slide_height - header_height - Inches(0.5)
    )
    text_frame = content_box.text_frame

    # Ajouter chaque bullet
    for i, bullet_text in enumerate(bullets):
        p = text_frame.add_paragraph()
        p.text = bullet_text
        p.level = 0  # niveau de puce
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0x00, 0x00, 0x00)  # noir
        if i == 0:
            # On peut mettre la première puce en semi-bold si on veut
            p.font.bold = False

def main():
    # Crée une présentation vierge
    prs = Presentation()

    # Paramètre optionnel: peut ajuster le format de la diapo (par défaut 13.333"x7.5")
    # Exemple format 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Construire les diapositives
    for slide_info in SLIDES_CONTENT:
        add_styled_slide(
            prs,
            title=slide_info["title"],
            bullets=slide_info["bullets"],
            theme_color=slide_info["theme_color"],
        )

    # Sauvegarder
    output_name = "ces2025_presentation_styled.pptx"
    prs.save(output_name)
    print(f"Fichier '{output_name}' généré avec succès.")


if __name__ == "__main__":
    main()

