#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
Script Python générant un PowerPoint (.pptx) avec python-pptx
incluant TOUT le contenu textuel donné (sans omission), réparti
sur 9 diapositives.
"""

from pptx import Presentation
from pptx.util import Inches, Pt

def add_bullet_slide(prs, title_text, bullet_list):
    """
    Ajoute une diapositive 'Title and Content' avec:
      - title_text : string affiché dans la zone de titre
      - bullet_list : liste de strings, chaque élément devient une puce
    """
    # Utilise la mise en page "Title and Content" (souvent index=1)
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # Titre
    slide.shapes.title.text = title_text

    # Zone de contenu pour les puces
    body_shape = slide.placeholders[1]
    text_frame = body_shape.text_frame

    # Si on a au moins une ligne, la première ira dans le paragraphe initial
    if bullet_list:
        text_frame.text = bullet_list[0]
    
    # Les puces suivantes vont dans de nouveaux paragraphes
    for bullet in bullet_list[1:]:
        p = text_frame.add_paragraph()
        p.text = bullet
        # p.level = 0 pour un niveau de puce standard

    # Ajustements optionnels (taille de police, etc.)
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(14)


def main():
    # 1) Créer une présentation vierge
    prs = Presentation()

    # -- SLIDE 1 : Titre + Sommaire --
    slide1_bullets = [
        "Introduction",
        "AI and Robotics in the Spotlight",
        "AI Robots for Social and Emotional Interaction",
        "AI Robots for Autonomy and Physical Assistance",
        "AI in Innovation and Creativity",
        "NVIDIA’s AI and Graphics Advancements",
        "Transition – Bridging AI Robotics and NVIDIA’s Tech",
        "Conclusion"
    ]
    add_bullet_slide(prs, "CES 2025: Innovations in AI and Robotics", slide1_bullets)

    # -- SLIDE 2 : Introduction – CES 2025 Overview --
    slide2_bullets = [
        (
            "Global Stage for Tech Innovation: CES 2025 in Las Vegas drew over 141,000 "
            "attendees and 4,500 exhibitors from around the world (ces.tech). As the "
            "premier consumer electronics show, it showcased breakthrough technologies "
            "across AI, robotics, smart homes, automotive, and more. Industry leaders "
            "and startups alike use CES to debut products and set trends, making it a "
            "bellwether event for the tech industry’s future."
        ),
        (
            "AI and Robotics in the Spotlight: A key theme of CES 2025 was the pervasive "
            "presence of artificial intelligence – “AI was everywhere,” turning once-"
            "speculative ideas into practical applications (rolandberger.com). From "
            "intelligent personal robots to autonomous vehicles, exhibitors demonstrated "
            "how advanced AI is enhancing devices’ capabilities. This year marked a shift "
            "toward AI-driven solutions addressing real-world challenges (e.g. aging "
            "populations, personalized entertainment), underlining CES’s role in shaping "
            "how technology meets human needs."
        ),
    ]
    add_bullet_slide(prs, "Introduction – CES 2025 Overview", slide2_bullets)

    # -- SLIDE 3 : AI and Robotics in the Spotlight --
    slide3_bullets = [
        (
            "This year, AI was everywhere at CES 2025, shifting towards real-world "
            "applications like aging population support and personalized entertainment. "
            "Robots, from personal companions to autonomous vehicles, highlighted how "
            "advanced AI capabilities are permeating consumer electronics."
        ),
    ]
    add_bullet_slide(prs, "AI and Robotics in the Spotlight", slide3_bullets)

    # -- SLIDE 4 : AI Robots for Social and Emotional Interaction --
    slide4_bullets = [
        (
            "Romi – A palm-sized emotional-support AI robot by Mixi (Japan) that engages "
            "in natural, empathic conversations. Designed to combat loneliness and anxiety, "
            "Romi can hold complex real-time dialogue without the user needing to pause – "
            "it even interjects with relevant comments based on what it “sees” and “hears” "
            "(ces.tech). Powered by in-house deep learning models, it strives to provide "
            "genuine empathy and companionship, using a cute animated face and playful "
            "gestures to simulate human-like interaction (reuters.com)."
        ),
        (
            "TCL “Ai Me” – A modular AI companion robot unveiled as a concept at CES 2025, "
            "aimed at supporting children, seniors, and families. This owl-eyed robot sits "
            "in a detachable “space capsule” base and can accompany users anywhere (even "
            "mounting in a smart car) (designboom.com). It has expressive digital eyes "
            "and a friendly voice, using cameras and sensors to recognize people and "
            "surroundings. Ai Me can record family moments, act as a home security "
            "sentry, and control smart home devices – a cute guardian and assistant "
            "fostering warmth and emotional connection (designboom.com)."
        ),
    ]
    add_bullet_slide(prs, "AI Robots for Social and Emotional Interaction", slide4_bullets)

    # -- SLIDE 5 : AI Robots for Autonomy and Physical Assistance --
    slide5_bullets = [
        (
            "Bosch “Revol” Smart Crib – An AI-powered connected cradle that helps parents "
            "care for infants. It monitors the baby’s vital signs (heart rate, breathing) "
            "using radar sensors and alerts caregivers via an app to any issues "
            "(homecrux.com). It can automatically rock or play lullabies when the baby "
            "cries and even detect changes in air quality or if the baby’s face is "
            "covered, notifying parents for intervention (housedigest.com). Bosch also "
            "designed the crib to be convertible into a height-adjustable desk as the "
            "child grows (housedigest.com)."
        ),
        (
            "Multi-Function Home Robots: CES 2025 showcased robots that serve as general "
            "home assistants. One prototype robot vacuum not only cleans floors but also "
            "acts as an air purifier, patrols the home for security, and even delivers "
            "items (it can literally bring you a sandwich on demand) (techradar.com). "
            "These all-in-one devices illustrate how AI and robotics are converging to "
            "autonomously manage household tasks – from cleaning and climate control to "
            "basic caregiving."
        ),
        (
            "Humanoid Helpers – “Mirokai” Robot: Enchanted Tools (France) debuted an "
            "updated Mirokai humanoid robot – a friendly, child-sized robot on wheels "
            "designed to assist and engage people. The bright orange Mirokai shown at "
            "CES has been deployed in hospitals as an aid for staff and companion for "
            "patients, and it can serve as a concierge in hospitality settings "
            "(interestingengineering.com, facebook.com). With an adorable, anime-inspired "
            "face and the ability to open doors, fetch objects, and interact socially, "
            "Mirokai provides both practical help and emotional comfort."
        ),
        (
            "Wearable Exoskeletons: AI-driven robotic exoskeleton suits offer new ways to "
            "assist those with mobility challenges. South Korea’s WIM exoskeleton (by "
            "WiRobotics) demonstrated a 78% improvement in walking speed and endurance "
            "among seniors (rockingrobots.com). These wearable robots use motors and "
            "sensors to augment the user’s strength and balance, hinting at a future "
            "where AI-powered exosuits restore mobility and independence."
        ),
    ]
    add_bullet_slide(prs, "AI Robots for Autonomy and Physical Assistance", slide5_bullets)

    # -- SLIDE 6 : AI in Innovation and Creativity --
    slide6_bullets = [
        (
            "AI-Generated Music: Startups and musicians at CES 2025 demonstrated AI "
            "systems composing music and soundscapes. Algorithms can generate original "
            "songs in real time based on a listener’s mood or a few sample inputs. "
            "Will.i.am highlighted the potential of AI-generated music, suggesting it "
            "could emotionally move people as deeply as legendary human artists "
            "(yahoo.com). These tools can aid creators by providing instant song drafts "
            "and personalize entertainment with mood-based DJ mixes."
        ),
        (
            "AI-Assisted Fashion Design: Generative AI design tools propose unique "
            "clothing patterns and styles, helping designers push beyond traditional "
            "boundaries (sganalytics.com). Concepts of virtual try-on experiences and AI "
            "stylists that recommend outfits were featured. For consumers, this means "
            "more personalization, and for creators, rapid prototyping of bold new "
            "designs. Fashion in 2025 embraces AI for creativity and efficiency."
        ),
        (
            "Interactive Holograms: Companies showcased 3D holographic displays that "
            "appear free-floating, no AR glasses required. UK-based Hypervsn wowed "
            "attendees with Blade Runner–style holograms – dancers, virtual assistants "
            "hovering mid-air (esquire.com). Future visions include holographic "
            "spokespeople, entertainers, and interactive AI avatars that you can gesture "
            "to. This may become a new medium for communication and art."
        ),
    ]
    add_bullet_slide(prs, "AI in Innovation and Creativity", slide6_bullets)

    # -- SLIDE 7 : NVIDIA’s AI and Graphics Advancements at CES 2025 --
    slide7_bullets = [
        (
            "GeForce RTX 5000 Series Unveiled: Jensen Huang presented next-gen GeForce "
            "RTX 50-series GPUs (Blackwell architecture). The RTX 5090 can do advanced "
            "path-traced lighting and neural rendering, delivering ultra-realistic "
            "graphics (esquire.com). The RTX 5070 at $550 performs on par with the prior "
            "$1,600 RTX 4090 (esquire.com), a 3× value improvement that makes high-end "
            "gaming and AI processing more accessible."
        ),
        (
            "DLSS 4 – Multi-Frame AI Upscaling: NVIDIA introduced DLSS 4 with Multi "
            "Frame Generation, where AI can generate up to 3 additional frames for every "
            "1 frame rendered (counterpointresearch.com). This yields massively higher "
            "frame rates and smoother motion. Together with improved ray tracing, "
            "AI-optimized shaders, and Reflex 2, DLSS 4 showcases how AI is "
            "revolutionizing real-time graphics."
        ),
        (
            "AI-Powered Robotics Platforms (NVIDIA Isaac & “Cosmos”): NVIDIA also "
            "revealed “Cosmos,” a suite of foundation AI models that generate "
            "photorealistic virtual worlds governed by physics (reuters.com). This "
            "synthetic data accelerates robotics training, reducing real-world trials. "
            "Cosmos is open-license, pairing with the NVIDIA Isaac platform for robot "
            "simulation and development. Such synergy between AI and graphics paves "
            "the way for faster innovation in robotics."
        ),
    ]
    add_bullet_slide(prs, "NVIDIA’s AI and Graphics Advancements at CES 2025", slide7_bullets)

    # -- SLIDE 8 : Transition – Bridging AI Robotics and NVIDIA’s Tech --
    slide8_bullets = [
        (
            "Synergy of Hardware and AI: The leaps in AI robotics are intertwined with "
            "advances in computing power from companies like NVIDIA. Many smart robots "
            "run on platforms enabled by AI chips and software, meaning improvements in "
            "GPUs and AI models boost robotic capabilities, while robotics’ real-time "
            "demands push chipmakers to innovate specialized hardware (rolandberger.com)."
        ),
        (
            "AI + Robotics Convergence: CES 2025 blurred the line between consumer tech "
            "and robotics. The same neural networks accelerating gaming also empower "
            "robot vision and speech. NVIDIA’s keynote launched gaming GPUs and AI "
            "training tools for robots in one go (reuters.com). As robots become common "
            "in homes and workplaces, they leverage better chips, refined algorithms, "
            "and cloud services, reinforcing a virtuous cycle of AI-driven innovation."
        ),
    ]
    add_bullet_slide(prs, "Transition – Bridging AI Robotics and NVIDIA’s Tech", slide8_bullets)

    # -- SLIDE 9 : Conclusion – Key Takeaways and Future Outlook --
    slide9_bullets = [
        (
            "Major Innovations Recap: CES 2025 showed how AI and robotics are "
            "transforming daily life. Emotional-support robots like Romi and Ai Me, "
            "autonomous devices (smart cribs, exoskeletons), AI-generated music, "
            "AI-assisted fashion, and NVIDIA’s powerful chips/platforms all point to "
            "a more connected, human-centric future."
        ),
        (
            "Benefits to People and Society: Social robots can alleviate loneliness and "
            "provide care for aging populations; assistive robots improve accessibility "
            "and safety; AI creative tools empower artists and personalize entertainment. "
            "Used thoughtfully, these technologies enhance quality of life, efficiency, "
            "and well-being."
        ),
        (
            "Challenges and Considerations: Privacy, reliability, and ethics must be "
            "addressed. Devices collecting personal data (cameras, mics) need secure "
            "handling. Society must ensure safe, trustworthy AI companions, navigate job "
            "displacement, and establish norms/regulations for AI behavior."
        ),
        (
            "Future Vision – “Smart” Everything: From more natural human–AI interaction "
            "to seamless home integration and advanced medical/rehab uses, the future "
            "hinted at by CES 2025 is one of synergy between people and AI-driven machines. "
            "By balancing innovation with responsible development, we may soon see "
            "intelligent robots and AI systems working hand-in-hand with humans, "
            "enriching lives in ways we’re just beginning to imagine."
        ),
    ]
    add_bullet_slide(prs, "Conclusion – Key Takeaways and Future Outlook", slide9_bullets)

    # 2) Enregistrer la présentation
    prs.save("ces2025_presentation.pptx")
    print("Fichier 'ces2025_presentation.pptx' généré avec succès.")

if __name__ == "__main__":
    main()

