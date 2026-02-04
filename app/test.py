from pptx import Presentation
from pptx.util import Inches

def add_slide(prs, title, content=None, bullet_points=None):
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]
    
    title_placeholder.text = title
    
    if content:
        content_placeholder.text = content
    
    if bullet_points:
        for point in bullet_points:
            p = content_placeholder.text_frame.add_paragraph()
            p.text = point

def generate_presentation():
    prs = Presentation()
    
    # Title Slide
    slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "AI and Robotics Innovations"
    slide.placeholders[1].text = "Exploring the Future of AI-Driven Robots and Graphics"
    
    # Section I: AI Robots for Social and Emotional Interaction
    add_slide(prs, "I. AI Robots for Social and Emotional Interaction")
    
    add_slide(prs, "1. Companion and Empathetic Robots")
    add_slide(prs, "Romi (Conversational Robot)", bullet_points=[
        "Engages in natural, fluid conversation.",
        "Detects emotions and responds accordingly.",
        "Helps reduce loneliness and stress."
    ])
    
    add_slide(prs, "AIME (Assistance Robot for Children and Seniors)", bullet_points=[
        "Interactive companion with emotional recognition.",
        "Monitors elderly people (medication reminders, health tracking).",
        "Can control smart home devices."
    ])
    
    # Section II: AI Robots for Autonomy and Physical Assistance
    add_slide(prs, "II. AI Robots for Autonomy and Physical Assistance")
    add_slide(prs, "1. Robots that Improve Daily Life")
    
    add_slide(prs, "Smart Connected Cradle", bullet_points=[
        "Advanced monitoring of babies (breathing, temperature, sleep).",
        "Automatically soothes the baby without parental intervention.",
        "Impact on modern parenting."
    ])
    
    add_slide(prs, "Multifunctional Smart Tower", bullet_points=[
        "Combines multiple devices in one (air purifier, speaker, light).",
        "AI adjusts functions based on the environment and user needs."
    ])
    
    add_slide(prs, "2. Exoskeletons and Physical Assistance Robots")
    add_slide(prs, "Eurotix (Lightweight Exoskeleton)", bullet_points=[
        "Muscle support and fatigue reduction.",
        "Applications in work, hiking, and rehabilitation."
    ])
    
    add_slide(prs, "Miroki (Humanoid Assistance Robot)", bullet_points=[
        "Developed by the French startup Enchanted Tools.",
        "1.23 meters, articulated arms, long ears, moves on a sphere.",
        "Carries loads up to 3 kg, performs repetitive or strenuous tasks.",
        "Eases the workload of healthcare personnel.",
        "A reassuring companion for patients, especially children in hospital settings."
    ])
    
    # Section III: AI Robots for Innovation and Creativity
    add_slide(prs, "III. AI Robots for Innovation and Creativity")
    add_slide(prs, "1. AI-Powered Creativity")
    
    add_slide(prs, "AI-Generated Music Piano Emage tomo", bullet_points=[
        "Composes music in real-time based on the atmosphere.",
        "Future possibilities of AI-assisted music composition."
    ])
    
    add_slide(prs, "AI-Generated Fashion Designs", bullet_points=[
        "Automated fashion design through artificial intelligence.",
        "Impact on the textile industry and clothing personalization."
    ])
    
    add_slide(prs, "2. AI Robots and Holograms")
    add_slide(prs, "HoloConnex (Interactive Holograms)", bullet_points=[
        "3D virtual presence for meetings, training, and entertainment.",
        "A step towards virtual teleportation."
    ])
    
    # Conclusion
    add_slide(prs, "Conclusion", bullet_points=[
        "Summary of major AI robotics advancements.",
        "Benefits: enhanced assistance, improved well-being, increased creativity.",
        "Challenges: dependence on machines, data privacy, balance between technology and human interaction.",
        "Future vision: AI robots will become increasingly integrated into daily life, but ethical challenges must be addressed."
    ])
    
    # Section IV: NVIDIA at CES 2025
    add_slide(prs, "IV. NVIDIA at CES 2025: AI, Graphics, and Robotics Innovations")
    
    add_slide(prs, "1. Next-Generation Graphics Technology")
    add_slide(prs, "GeForce RTX 5000 Series", bullet_points=[
        "Advanced Performance: Features up to 21,760 CUDA cores and 32 GB of GDDR7 memory for high-speed, high-performance processing.",
        "Enhanced Visuals: Delivers superior gaming and creative content experiences, particularly in 4K with advanced ray tracing.",
        "Power Efficiency: Optimized for better energy consumption while maximizing output."
    ])
    
    add_slide(prs, "2. DLSS 4: Revolutionizing AI-Powered Graphics")
    add_slide(prs, "Overview of DLSS 4", bullet_points=[
        "Utilizes deep learning models to upscale lower-resolution images, enhancing performance without compromising visual fidelity."
    ])
    
    add_slide(prs, "Key Features of DLSS 4", bullet_points=[
        "Multi Frame Generation (MFG): Generates up to three additional frames for every traditionally rendered frame, increasing frame rates significantly (up to 8x improvement) for ultra-smooth gameplay.",
        "Transformer-Based AI Models: Employs sophisticated AI models for more accurate image reconstruction and reduced visual artifacts.",
        "Ray Reconstruction: Enhances the quality and realism of ray-traced lighting and reflections for a more immersive gaming experience."
    ])
    
    add_slide(prs, "Impact on Gaming and Creative Industries")
    
    add_slide(prs, "3. NVIDIA and AI in Robotics")
    add_slide(prs, "AI-Powered Robotics Platforms", bullet_points=[
        "NVIDIA Isaac: A comprehensive robotics development platform providing simulation and real-time AI processing for autonomous machines."
    ])
    
    add_slide(prs, "Applications in Robotics", bullet_points=[
        "Automation and Adaptability: Robots can learn from their environments and adapt to complex tasks.",
        "Healthcare and Logistics: AI-powered robots assist in hospitals, factories, and warehouses, improving efficiency and safety."
    ])
    
    add_slide(prs, "4. Conclusion", bullet_points=[
        "NVIDIA’s Vision: Reinforces their position as a leader in AI-driven innovation for gaming, creative industries, and robotics.",
        "Future Impact: Demonstrates how AI will continue to reshape the landscape of entertainment and technology.",
        "Ethical Considerations: Emphasizes the need to balance technological advancement with ethical responsibility, particularly regarding data privacy and AI dependency."
    ])
    
    prs.save("AI_Robotics_Presentation.pptx")
    print("Presentation saved as AI_Robotics_Presentation.pptx")

# Run the script
generate_presentation()

